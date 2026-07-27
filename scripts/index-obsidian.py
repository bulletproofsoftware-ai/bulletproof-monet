#!/usr/bin/env python3
"""
Obsidian Vault Indexer for RAG (c2 server version)
Scans markdown, PDF, and MS Office files, generates embeddings via Ollama, and stores in Qdrant.
"""

import os
import sys
import json
import hashlib
import requests
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional, Tuple
import uuid

# Configuration
OBSIDIAN_VAULT = os.environ.get("OBSIDIAN_VAULT", "/root/obsidian-vault")
QDRANT_URL = os.environ.get("QDRANT_URL", "http://localhost:6334")
QDRANT_API_KEY = os.environ.get("QDRANT_API_KEY", "")
OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://localhost:11434")
COLLECTION_NAME = "obsidian_docs"
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
MAX_FILE_SIZE = 10 * 1024 * 1024  # Skip files > 10MB (image-heavy PDFs)

SUPPORTED_EXTENSIONS = {".md", ".pdf", ".docx", ".xlsx", ".pptx"}
SKIP_DIRS = {".obsidian", ".git", ".trash", "node_modules", ".claude", "_Templates", "copilot"}
SKIP_PATTERNS = {".tmp.driveupload"}


def get_file_hash(content: str) -> str:
    return hashlib.sha256(content.encode()).hexdigest()


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    if len(text) <= chunk_size:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if end < len(text):
            para_break = chunk.rfind("\n\n")
            if para_break > chunk_size // 2:
                chunk = chunk[:para_break]
                end = start + para_break
            else:
                for sep in [". ", "! ", "? ", "\n"]:
                    sent_break = chunk.rfind(sep)
                    if sent_break > chunk_size // 2:
                        chunk = chunk[:sent_break + len(sep)]
                        end = start + sent_break + len(sep)
                        break
        chunks.append(chunk.strip())
        start = end - overlap
    return [c for c in chunks if c]


def extract_text_from_pdf(file_path: Path) -> Tuple[str, str]:
    try:
        import pdfplumber
        import signal
        def _timeout_handler(signum, frame):
            raise TimeoutError("PDF extraction timed out")
        text_parts = []
        old_handler = signal.signal(signal.SIGALRM, _timeout_handler)
        signal.alarm(60)  # 60s timeout per PDF
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
        finally:
            signal.alarm(0)
            signal.signal(signal.SIGALRM, old_handler)
        content = "\n\n".join(text_parts)
        title = file_path.stem
        if content:
            first_line = content.split("\n")[0].strip()
            if first_line and len(first_line) < 200:
                title = first_line
        return content, title
    except TimeoutError:
        print(f"  Skipped (PDF extraction timeout)")
        return "", file_path.stem
    except Exception as e:
        print(f"  Error extracting PDF: {e}")
        return "", file_path.stem


def extract_text_from_docx(file_path: Path) -> Tuple[str, str]:
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        content = "\n\n".join(paragraphs)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    content += "\n" + row_text
        title = file_path.stem
        if paragraphs:
            first_para = paragraphs[0].strip()
            if first_para and len(first_para) < 200:
                title = first_para
        return content, title
    except Exception as e:
        print(f"  Error extracting DOCX: {e}")
        return "", file_path.stem


def extract_text_from_xlsx(file_path: Path) -> Tuple[str, str]:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"## Sheet: {sheet_name}")
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    text_parts.append(row_text)
        content = "\n".join(text_parts)
        return content, file_path.stem
    except Exception as e:
        print(f"  Error extracting XLSX: {e}")
        return "", file_path.stem


def extract_text_from_pptx(file_path: Path) -> Tuple[str, str]:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        title = file_path.stem
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                text_parts.append(f"## Slide {i}")
                text_parts.extend(slide_texts)
            if i == 1 and slide_texts:
                potential_title = slide_texts[0]
                if len(potential_title) < 200:
                    title = potential_title
        content = "\n\n".join(text_parts)
        return content, title
    except Exception as e:
        print(f"  Error extracting PPTX: {e}")
        return "", file_path.stem


def extract_text_from_markdown(file_path: Path) -> Tuple[str, str]:
    try:
        content = file_path.read_text(encoding="utf-8")
        title = file_path.stem
        for line in content.split("\n"):
            if line.startswith("# "):
                title = line[2:].strip()
                break
        return content, title
    except Exception as e:
        print(f"  Error reading markdown: {e}")
        return "", file_path.stem


def extract_text(file_path: Path) -> Tuple[str, str]:
    ext = file_path.suffix.lower()
    extractors = {
        ".md": extract_text_from_markdown,
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".xlsx": extract_text_from_xlsx,
        ".pptx": extract_text_from_pptx,
    }
    return extractors.get(ext, lambda p: ("", p.stem))(file_path)


def generate_embedding(text: str) -> Optional[List[float]]:
    try:
        response = requests.post(
            f"{OLLAMA_URL}/api/embeddings",
            json={"model": "nomic-embed-text", "prompt": text},
            timeout=30
        )
        response.raise_for_status()
        return response.json()["embedding"]
    except Exception as e:
        print(f"  Error generating embedding: {e}")
        return None


def store_in_qdrant(points: List[Dict]) -> bool:
    try:
        response = requests.put(
            f"{QDRANT_URL}/collections/{COLLECTION_NAME}/points",
            headers={"api-key": QDRANT_API_KEY, "Content-Type": "application/json"},
            json={"points": points},
            timeout=30
        )
        response.raise_for_status()
        return True
    except Exception as e:
        print(f"  Error storing in Qdrant: {e}")
        return False


def delete_file_points(file_path: str) -> bool:
    try:
        response = requests.post(
            f"{QDRANT_URL}/collections/{COLLECTION_NAME}/points/delete",
            headers={"api-key": QDRANT_API_KEY, "Content-Type": "application/json"},
            json={"filter": {"must": [{"key": "file_path", "match": {"value": file_path}}]}},
            timeout=30
        )
        response.raise_for_status()
        return True
    except Exception as e:
        print(f"  Error deleting old points: {e}")
        return False


def get_existing_hashes() -> Dict[str, str]:
    hashes = {}
    try:
        offset = None
        while True:
            payload = {
                "limit": 100,
                "with_payload": {"include": ["file_path", "content_hash"]},
                "with_vector": False
            }
            if offset:
                payload["offset"] = offset
            response = requests.post(
                f"{QDRANT_URL}/collections/{COLLECTION_NAME}/points/scroll",
                headers={"api-key": QDRANT_API_KEY, "Content-Type": "application/json"},
                json=payload, timeout=30
            )
            response.raise_for_status()
            result = response.json()["result"]
            for point in result.get("points", []):
                p = point.get("payload", {})
                fp = p.get("file_path")
                ch = p.get("content_hash")
                if fp and ch:
                    hashes[fp] = ch
            offset = result.get("next_page_offset")
            if not offset:
                break
    except Exception as e:
        print(f"Warning: Could not get existing hashes: {e}")
    return hashes


def should_skip_path(path: Path) -> bool:
    for part in path.parts:
        if part in SKIP_DIRS:
            return True
        for pattern in SKIP_PATTERNS:
            if pattern in part:
                return True
    return False


def index_file(file_path: Path, relative_path: str, existing_hashes: Dict[str, str]) -> int:
    # Skip files that are too large (usually image-heavy PDFs)
    try:
        file_size = file_path.stat().st_size
        if file_size > MAX_FILE_SIZE:
            print(f"  Skipped (>{MAX_FILE_SIZE//1024//1024}MB): {file_size//1024//1024}MB")
            return 0
    except OSError:
        pass
    content, title = extract_text(file_path)
    if not content.strip():
        return 0
    content_hash = get_file_hash(content)
    if relative_path in existing_hashes and existing_hashes[relative_path] == content_hash:
        return 0
    delete_file_points(relative_path)
    chunks = chunk_text(content)
    points = []
    for i, chunk in enumerate(chunks):
        embedding = generate_embedding(chunk)
        if embedding:
            points.append({
                "id": str(uuid.uuid4()),
                "vector": embedding,
                "payload": {
                    "file_path": relative_path,
                    "title": title,
                    "content": chunk,
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "content_hash": content_hash,
                    "indexed_at": datetime.now().isoformat(),
                    "file_size": len(content),
                    "file_type": file_path.suffix.lower()
                }
            })
    if points and store_in_qdrant(points):
        return len(points)
    return 0


def find_supported_files(vault_path: Path) -> List[Path]:
    files = []
    for ext in SUPPORTED_EXTENSIONS:
        for file_path in vault_path.rglob(f"*{ext}"):
            if not should_skip_path(file_path):
                files.append(file_path)
    return sorted(files, key=lambda p: str(p))


def main():
    if not QDRANT_API_KEY:
        print("ERROR: QDRANT_API_KEY not set")
        sys.exit(1)

    print(f"Obsidian Vault Indexer (c2)")
    print(f"=" * 50)
    print(f"Vault: {OBSIDIAN_VAULT}")
    print(f"Qdrant: {QDRANT_URL}")
    fmts = ", ".join(SUPPORTED_EXTENSIONS); print(f"Supported formats: {fmts}")
    print()

    print("Checking existing index...")
    existing_hashes = get_existing_hashes()
    print(f"Found {len(existing_hashes)} previously indexed file chunks")
    print()

    vault_path = Path(OBSIDIAN_VAULT)
    files = find_supported_files(vault_path)

    type_counts = {}
    for f in files:
        ext = f.suffix.lower()
        type_counts[ext] = type_counts.get(ext, 0) + 1

    print(f"Found {len(files)} files to process:")
    for ext, count in sorted(type_counts.items()):
        print(f"  {ext}: {count}")
    print()

    total_chunks = 0
    indexed_files = 0
    skipped_files = 0
    error_files = 0

    for i, file_path in enumerate(files, 1):
        relative_path = str(file_path.relative_to(vault_path))
        ext = file_path.suffix.lower()
        display = relative_path[:60]
        print(f"[{i}/{len(files)}] [{ext}] {display}...")

        try:
            chunks = index_file(file_path, relative_path, existing_hashes)
            if chunks > 0:
                print(f"  Indexed {chunks} chunks")
                total_chunks += chunks
                indexed_files += 1
            else:
                skipped_files += 1
        except Exception as e:
            print(f"  Error: {e}")
            error_files += 1

    print()
    print(f"=" * 50)
    print(f"Indexing complete!")
    print(f"  Files indexed: {indexed_files}")
    print(f"  Files skipped (unchanged): {skipped_files}")
    print(f"  Files with errors: {error_files}")
    print(f"  Total new chunks: {total_chunks}")


if __name__ == "__main__":
    main()
